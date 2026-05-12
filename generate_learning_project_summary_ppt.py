from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x18, 0x18, 0x1B)
GRAY = RGBColor(0x5F, 0x66, 0x73)
LIGHT_BG = RGBColor(0xF6, 0xF7, 0xFB)
BLUE = RGBColor(0x25, 0x63, 0xEB)
BLUE_DARK = RGBColor(0x1E, 0x40, 0xAF)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
RED = RGBColor(0xDC, 0x26, 0x26)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
PALE_BLUE = RGBColor(0xDB, 0xEA, 0xFE)


def set_bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, text, size=16, color=BLACK, bold=False,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    return shape


def add_bullets(slide, x, y, w, h, items, size=13, color=BLACK, gap=5):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(gap)
    return shape


def title(slide, text, subtitle=None, accent=BLUE):
    textbox(slide, 0.62, 0.34, 12.0, 0.52, text, size=26, color=accent, bold=True)
    if subtitle:
        textbox(slide, 0.65, 0.92, 11.7, 0.36, subtitle, size=11.5, color=GRAY)


def card(slide, x, y, w, h, heading, body, accent=BLUE, heading_size=15, body_size=11.5):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = CARD
    rect.line.color.rgb = BORDER
    rect.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    textbox(slide, x + 0.22, y + 0.15, w - 0.34, 0.33, heading, size=heading_size, color=accent, bold=True)
    add_bullets(slide, x + 0.22, y + 0.56, w - 0.36, h - 0.64, body, size=body_size, color=BLACK, gap=3)
    return rect


def tag(slide, x, y, text, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.0), Inches(0.28))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    textbox(slide, x, y + 0.03, 1.0, 0.18, text, size=8.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def section_label(slide, text):
    textbox(slide, 0.65, 6.98, 5.8, 0.22, text, size=8.5, color=GRAY)


def add_two_column(slide, left_title, left_items, right_title, right_items, left_color=BLUE, right_color=PURPLE):
    card(slide, 0.75, 1.38, 5.75, 5.45, left_title, left_items, left_color, body_size=12)
    card(slide, 6.82, 1.38, 5.75, 5.45, right_title, right_items, right_color, body_size=12)


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLUE)
    textbox(slide, 0.8, 1.28, 11.8, 0.75, "大模型应用开发", size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.8, 2.18, 11.8, 0.58, "4 周学习项目详细复盘", size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 1.25, 3.08, 10.8, 0.88,
            "Python API · Prompt · 结构化输出 · Eval · RAG · FastAPI · Agent · LangGraph · MCP · 部署",
            size=15.5, color=PALE_BLUE, align=PP_ALIGN.CENTER)
    textbox(slide, 1.35, 5.15, 10.7, 0.35, "学习仓库: D:\\code\\py", size=13, color=PALE_BLUE, align=PP_ALIGN.CENTER)
    textbox(slide, 1.35, 5.62, 10.7, 0.35, "用途: 复习路线、重点概念、项目产出、后续打磨方向", size=13, color=PALE_BLUE, align=PP_ALIGN.CENTER)


def how_to_use(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "这份 PPT 怎么看", "不是展示用的空泛总结，而是学习复盘用的路线图。")
    add_two_column(
        slide,
        "每个阶段看 4 件事",
        [
            "学了什么: 每周主题和每天内容",
            "重点是什么: 哪些概念必须真正理解",
            "代码在哪里: 对应目录和项目文件",
            "怎么自检: 能否解释、运行、评估和改进",
            "",
            "建议看法:",
            "先看路线总览 → 再按 Week 逐页复盘 → 最后看下一阶段任务。",
        ],
        "掌握标准",
        [
            "不是“看过课程文档”",
            "不是“跑过一次 demo”",
            "而是能做到:",
            "1. 不看答案能复现最小实现",
            "2. 能讲清技术取舍",
            "3. 能定位失败原因",
            "4. 能把 demo 包装成作品",
            "5. 能解释成本、延迟、稳定性问题",
        ],
        BLUE,
        GREEN,
    )
    section_label(slide, "How to use")


def agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "目录", "从路线判断到每周学习重点，再到作品集和下一步计划。")
    rows = [
        ("01", "路线总览", "4 周为什么这样排，核心能力链路是什么"),
        ("02", "Week 1", "Python 基础、原生 API、streaming、CLI 工具"),
        ("03", "Week 2", "Prompt、结构化输出、Tool Use、Pydantic、Eval、LangChain"),
        ("04", "Week 3", "RAG、向量库、FastAPI、SSE、Docs Copilot"),
        ("05", "Week 4", "Agent、ReAct、LangGraph、MCP、部署、Workflow Assistant"),
        ("06", "横向能力", "评估体系、作品集、当前问题、下一阶段计划"),
    ]
    y = 1.45
    for num, head, desc in rows:
        textbox(slide, 0.95, y, 0.7, 0.32, num, size=17, color=BLUE, bold=True)
        textbox(slide, 1.8, y, 2.65, 0.32, head, size=17, color=BLACK, bold=True)
        textbox(slide, 4.45, y + 0.02, 7.6, 0.30, desc, size=13.2, color=GRAY)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(y + 0.48), Inches(11.1), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER
        line.line.fill.background()
        y += 0.82
    section_label(slide, "Agenda")


def route_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "4 周路线总览", "这条路线的核心是: 原生能力打底, 评估贯穿, 每周沉淀作品。")
    cards = [
        ("Week 1\nPython + API", ["Python 基础", "环境变量", "JSON/文件 IO", "原生模型调用", "streaming"], BLUE),
        ("Week 2\nPrompt + Eval", ["Prompt 模板", "结构化输出", "Tool Use", "Pydantic", "回归评估"], PURPLE),
        ("Week 3\nRAG + Web", ["chunk/embedding", "Chroma 检索", "FastAPI", "SSE 前端", "Docs Copilot"], GREEN),
        ("Week 4\nAgent + Deploy", ["ReAct", "LangGraph", "MCP", "项目评估", "部署上线"], ORANGE),
    ]
    x = 0.75
    for head, items, color in cards:
        card(slide, x, 1.38, 2.82, 4.75, head, items, color, body_size=12)
        x += 3.08
    textbox(slide, 0.9, 6.35, 11.7, 0.35, "能力链路: 输入设计 → 模型调用 → 结构化输出 → 检索/工具 → 应用接口 → 评估优化 → 部署展示", size=14, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    section_label(slide, "Route")


def route_principles(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "路线判断: 为什么这样学", "如果顺序错了，很容易学成“框架 API 背诵”，而不是 AI 应用开发。")
    add_two_column(
        slide,
        "正确顺序",
        [
            "1. 先原生 SDK",
            "理解 messages、streaming、token、错误处理。",
            "",
            "2. 再学结构化输出",
            "让模型输出可解析、可校验、可落库。",
            "",
            "3. 评估前置",
            "用固定输入检查 prompt 是否退化。",
            "",
            "4. 做完整项目",
            "输入、后端、前端、评估、部署都要闭环。",
        ],
        "容易走偏的地方",
        [
            "一上来学 3 个 Agent 框架",
            "只看 RAG 概念，不做失败诊断",
            "只跑 demo，不写 eval_cases",
            "只写代码，不整理 README",
            "没有环境说明，换机器跑不起来",
            "不会解释延迟、成本、稳定性",
            "",
            "结论: 先能交付，再扩展框架视野。",
        ],
        BLUE,
        RED,
    )
    section_label(slide, "Learning principles")


def week_overview(prs, week, heading, goal, days, output, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, f"Week {week}: {heading}", goal, accent)
    card(slide, 0.75, 1.35, 5.9, 5.35, "每日学习内容", days, accent, body_size=12)
    card(slide, 6.95, 1.35, 5.65, 5.35, "本周产出", output, GREEN, body_size=12)
    section_label(slide, f"Week {week} overview")


def week1_overview(prs):
    week_overview(
        prs, 1, "Python 基础 + LLM 初体验",
        "目标: 补齐 Python 必需基础，并独立完成一次模型 API 调用和流式输出。",
        [
            "Day 1: 环境、venv、.env、第一次 API 调用",
            "Day 2: str/list/dict/function 等基础语法",
            "Day 3: 文件 IO、JSON、异常处理",
            "Day 4: 模块、类、包结构",
            "Day 5: async/await、并发、streaming",
            "Day 6-7: AI 笔记生成器项目",
        ],
        [
            "最小模型调用脚本",
            "命令行聊天机器人",
            "文件扫描 + AI 分析 + JSON/Markdown 输出",
            "模块化项目: chatbot/",
            "练习目录: week1-practice",
            "课程目录: week1-python-ai-basics",
        ],
        BLUE,
    )


def week1_details(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Week 1 重点: Python 只学 AI 应用必需部分", "不要把时间花在完整 Python 教程上，先覆盖项目开发高频能力。", BLUE)
    card(slide, 0.75, 1.35, 3.8, 5.2, "必须掌握", [
        "Path / open / with 文件读写",
        "json.load / loads / dump / dumps",
        "try / except / finally",
        "list/dict 遍历和格式化",
        "class 封装可复用逻辑",
        "import 与包结构",
        "asyncio.gather 并发调用",
    ], BLUE)
    card(slide, 4.75, 1.35, 3.8, 5.2, "LLM API 关键点", [
        "client 初始化和 API Key 管理",
        "model / messages / max_tokens",
        "system prompt 与 user prompt",
        "streaming 返回不是一次性文本",
        "错误处理和超时重试",
        "把模型输出保存为 JSON/Markdown",
    ], PURPLE)
    card(slide, 8.75, 1.35, 3.8, 5.2, "自检问题", [
        "能否独立创建 venv?",
        "能否解释 .env 为什么不能提交?",
        "能否把一个目录下的 .py 文件扫出来?",
        "能否并发调用多个文件分析?",
        "能否把 AI 结果保存成结构化文件?",
    ], ORANGE)
    section_label(slide, "Week 1 details")


def week1_project(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Week 1 项目: AI 笔记生成器", "第一个作品不是聊天机器人，而是能处理真实文件的命令行工具。", BLUE)
    add_two_column(
        slide,
        "项目流程",
        [
            "1. 扫描指定目录",
            "2. 过滤 .py / .md / .txt 等文件",
            "3. 读取文件内容并截断过长文本",
            "4. 调用模型生成摘要/重点",
            "5. 保存 notes_data.json",
            "6. 生成 notes_report.md",
            "",
            "训练点: 文件 IO、JSON、API、异常处理、模块拆分。",
        ],
        "代码定位",
        [
            "week1-practice/note_generator/",
            "  scanner.py: 扫描文件",
            "  generator.py: 生成笔记",
            "  config.py: 配置项",
            "",
            "week1-practice/output/",
            "  notes_data.json",
            "  notes_report.md",
            "  week1_summary.pptx",
            "",
            "后续可改进: 加进度条、并发、失败重试。",
        ],
        BLUE,
        GREEN,
    )
    section_label(slide, "Week 1 project")


def week2_overview(prs):
    week_overview(
        prs, 2, "Prompt Engineering + 结构化输出 + 评估",
        "目标: 从会聊天升级到能稳定抽取、校验、评估和调试。",
        [
            "Day 1: Few-shot、CoT、System Prompt",
            "Day 2: 结构化输出、Tool Use、Pydantic",
            "Day 3: eval_cases、日志、失败样本",
            "Day 4: LangChain Chain、PromptTemplate、Parser",
            "Day 5: LangChain Agent、自定义工具",
            "Day 6-7: AI 数据助手 / 简历分析助手",
        ],
        [
            "商品信息提取器",
            "简历抽取器",
            "智能客服路由器",
            "Prompt 评估脚本",
            "week2-project: 简历分析 Agent",
            "日志文件: api_log_*.jsonl",
        ],
        PURPLE,
    )


def week2_prompt(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Week 2 重点: Prompt 要工程化", "Prompt 不是写一段自然语言，而是给模型定义任务、边界和可检查输出。", PURPLE)
    card(slide, 0.75, 1.35, 3.8, 5.2, "Prompt 组成", [
        "角色: 你是谁",
        "目标: 要完成什么任务",
        "上下文: 输入材料是什么",
        "约束: 不能做什么",
        "格式: 必须怎么输出",
        "示例: Few-shot 对齐风格",
        "失败处理: 不确定时怎么回答",
    ], PURPLE)
    card(slide, 4.75, 1.35, 3.8, 5.2, "结构化输出", [
        "Prompt 约束 JSON: 简单但不稳",
        "Tool Use / Function Calling: 更适合生产",
        "Pydantic: 对字段做类型校验",
        "clean_markdown: 清理 ```json 代码块",
        "json.loads: 解析失败要有 fallback",
        "字段缺失要返回明确错误",
    ], BLUE)
    card(slide, 8.75, 1.35, 3.8, 5.2, "常见坑", [
        "只要求“输出 JSON”但不给 schema",
        "让模型自由发挥字段名",
        "解析失败直接崩溃",
        "没有保存原始响应",
        "没有固定测试用例",
        "prompt 改坏了没人知道",
    ], RED)
    section_label(slide, "Week 2 prompt")


def week2_eval_langchain(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Week 2 重点: 评估 + LangChain", "LangChain 是工具；评估才是让系统持续变好的机制。", PURPLE)
    add_two_column(
        slide,
        "评估应该怎么做",
        [
            "准备 eval_cases.json",
            "每条包含 input、expected_keywords、tags",
            "运行脚本批量调用模型",
            "检查关键词命中率",
            "记录 latency 和失败样本",
            "每次改 prompt 后跑回归",
            "",
            "目标: 让优化有数据, 不靠感觉。",
        ],
        "LangChain 学什么",
        [
            "PromptTemplate: 模板变量",
            "OutputParser: 输出解析",
            "LCEL: prompt | model | parser",
            "Memory: 对话上下文",
            "Tool: 把函数暴露给模型",
            "Agent: 模型决定调用哪个工具",
            "",
            "注意: 先理解概念, 不要被框架牵着走。",
        ],
        GREEN,
        BLUE,
    )
    section_label(slide, "Week 2 eval + LangChain")


def week2_project(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Week 2 项目: AI 简历分析助手", "把结构化抽取、工具调用、评估和日志组合成一个完整小项目。", PURPLE)
    card(slide, 0.75, 1.35, 3.8, 5.2, "功能", [
        "输入简历文本",
        "提取姓名、年限、技能、学历、薪资",
        "估算薪资范围",
        "评估经验级别",
        "输出完整分析报告",
        "verbose 模式展示 Agent 步骤",
    ], PURPLE)
    card(slide, 4.75, 1.35, 3.8, 5.2, "代码定位", [
        "week2-project/main.py: CLI 入口",
        "week2-project/agent.py: LangGraph ReAct Agent",
        "week2-project/tools.py: 工具集合",
        "week2-project/extractor.py: 抽取逻辑",
        "week2-project/eval_runner.py: 评估脚本",
        "week2-project/logger.py: 日志记录",
    ], BLUE)
    card(slide, 8.75, 1.35, 3.8, 5.2, "下一步优化", [
        "删除调试输出 promptprompt",
        "补更严格 Pydantic schema",
        "拆分 CLI 与核心服务",
        "把日志字段标准化",
        "补 README 运行步骤",
        "增加失败案例分析",
    ], ORANGE)
    section_label(slide, "Week 2 project")


def week3_overview(prs):
    week_overview(
        prs, 3, "RAG + FastAPI + 全栈整合",
        "目标: 做出能检索私有文档、显示引用来源、支持前端流式交互的 AI 应用。",
        [
            "Day 1: RAG 基础、chunk、embedding",
            "Day 2: Chroma / FAISS 最小 Pipeline",
            "Day 3: top-k、阈值、引用、Rerank 概念",
            "Day 4: FastAPI 路由、Pydantic、CORS",
            "Day 5: 前端 SSE 流式接收",
            "Day 6-7: Docs Copilot",
        ],
        [
            "RAGPipeline 类封装",
            "文档索引脚本 indexer.py",
            "FastAPI 问答接口",
            "SSE 流式聊天页面",
            "引用来源展示",
            "RAG 评估用例",
        ],
        GREEN,
    )


def week3_rag(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Week 3 重点: RAG 不是背概念", "真正重要的是知道回答差时，问题出在检索、上下文还是生成。", GREEN)
    card(slide, 0.75, 1.35, 3.8, 5.2, "RAG 流程", [
        "文档加载",
        "分块 chunking",
        "embedding 向量化",
        "向量库存储",
        "用户问题向量化",
        "top-k 检索",
        "拼接 context",
        "LLM 基于文档回答",
    ], GREEN)
    card(slide, 4.75, 1.35, 3.8, 5.2, "调参点", [
        "chunk_size: 太大噪音多, 太小语义断",
        "overlap: 防止关键信息被切断",
        "top_k: 太少不完整, 太多噪音多",
        "threshold: 不相关时拒答",
        "rerank: 先粗召回, 再精排",
        "prompt: 明确只能基于文档回答",
    ], PURPLE)
    card(slide, 8.75, 1.35, 3.8, 5.2, "失败诊断", [
        "没搜到: chunk/embedding/query 问题",
        "搜错了: top_k/阈值/rerank 问题",
        "搜到了但答错: prompt 问题",
        "答得太长: max_tokens/格式约束",
        "编造: 缺拒答规则和引用约束",
        "慢: 检索数、上下文长度、模型选择",
    ], RED)
    section_label(slide, "Week 3 RAG")


def week3_fastapi_frontend(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Week 3 重点: FastAPI + SSE 前端", "这一步把 AI 能力变成用户可用的 Web 应用。", GREEN)
    add_two_column(
        slide,
        "FastAPI 要点",
        [
            "app = FastAPI()",
            "@app.get / @app.post 路由",
            "Pydantic 请求体和响应模型",
            "CORS 跨域配置",
            "HTMLResponse 返回页面",
            "EventSourceResponse / StreamingResponse",
            "自动生成 /docs 接口文档",
            "",
            "Express 类比: 路由 + middleware + JSON response。",
        ],
        "前端要点",
        [
            "fetch 发送问题",
            "ReadableStream 读取流",
            "TextDecoder 解码 token",
            "逐字追加到页面",
            "显示 loading/error/empty 状态",
            "展示 sources 引用来源",
            "可加 AbortController 停止生成",
            "",
            "前端优势: AI 产品体验主要靠状态管理和流式 UI。",
        ],
        BLUE,
        GREEN,
    )
    section_label(slide, "Week 3 FastAPI + Frontend")


def week3_project(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Week 3 项目: Docs Copilot", "这是第一个完整全栈 AI 应用: 文档 → 检索 → 回答 → 引用 → 前端。", GREEN)
    card(slide, 0.75, 1.35, 3.8, 5.2, "项目结构", [
        "week3-project/rag.py",
        "week3-project/indexer.py",
        "week3-project/main.py",
        "week3-project/models.py",
        "week3-project/templates/index.html",
        "week3-project/docs/",
        "week3-project/eval_runner.py",
    ], GREEN)
    card(slide, 4.75, 1.35, 3.8, 5.2, "核心功能", [
        "索引 Markdown 文档",
        "Chroma collection 存储",
        "语义检索 top-k",
        "基于 context 回答",
        "引用来源展示",
        "SSE 流式输出",
        "评估检索和生成质量",
    ], BLUE)
    card(slide, 8.75, 1.35, 3.8, 5.2, "改进方向", [
        "当前 Chroma 是内存模式",
        "可改 PersistentClient",
        "不要启动时 delete_collection",
        "补上传文档 UI",
        "增加真实知识库评估集",
        "加入阈值拒答和来源置信度",
    ], ORANGE)
    section_label(slide, "Week 3 project")


def week4_overview(prs):
    week_overview(
        prs, 4, "Agent + LangGraph + MCP + 部署",
        "目标: 理解 Agent 边界，用 LangGraph 写可控工作流，并具备部署和项目评估意识。",
        [
            "Day 1: Agent Loop、ReAct、停止条件",
            "Day 2: LangGraph State/Node/Edge",
            "Day 3: MCP 协议和工具标准化",
            "Day 4: 项目级评估与优化",
            "Day 5: 部署、环境变量、日志",
            "Day 6-7: AI Workflow Assistant",
        ],
        [
            "手写 SimpleAgent",
            "LangGraph 工作流 demo",
            "MCP 实验",
            "eval_runner.py 项目评估",
            "FastAPI + SSE workflow assistant",
            "部署材料和 README",
        ],
        ORANGE,
    )


def week4_agent(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Week 4 重点: Agent 要知道边界", "Agent 不是更高级的万能调用方式，它只是“循环 + 判断 + 工具”。", ORANGE)
    card(slide, 0.75, 1.35, 3.8, 5.2, "ReAct 流程", [
        "Thought: 分析当前情况",
        "Action: 选择工具",
        "Action Input: 工具参数",
        "Observation: 工具返回结果",
        "循环直到 Final Answer",
        "",
        "核心: prompt 格式 + parser + tools registry。",
    ], ORANGE)
    card(slide, 4.75, 1.35, 3.8, 5.2, "适合 Agent", [
        "多步骤研究任务",
        "路径不确定",
        "需要多个工具组合",
        "需要根据中间结果决定下一步",
        "例如: 查资料 → 计算 → 汇总",
    ], GREEN)
    card(slide, 8.75, 1.35, 3.8, 5.2, "不适合 Agent", [
        "简单翻译/摘要",
        "固定格式抽取",
        "对延迟极敏感",
        "工具很少或步骤固定",
        "成本敏感场景",
        "",
        "这些用 Chain 或一次调用更稳。",
    ], RED)
    section_label(slide, "Week 4 Agent")


def week4_langgraph_mcp(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Week 4 重点: LangGraph + MCP", "LangGraph 解决可控编排，MCP 解决工具接入标准化。", ORANGE)
    add_two_column(
        slide,
        "LangGraph 学习重点",
        [
            "State: 全局状态对象",
            "Node: 处理状态的函数",
            "Edge: 节点连接",
            "Conditional Edge: 条件路由",
            "END: 结束节点",
            "checkpointer: 保存状态",
            "interrupt_before: 人工介入",
            "",
            "适合: 生产级、可观测、可暂停的 Agent 工作流。",
        ],
        "MCP 学习重点",
        [
            "MCP = Model Context Protocol",
            "Client: AI 应用",
            "Server: 工具/数据提供方",
            "tools/list: 自动发现工具",
            "tools/call: 调用工具",
            "比普通 function calling 更标准化",
            "",
            "当前阶段: 了解定位、跑通实验即可，不要过度深挖。",
        ],
        BLUE,
        PURPLE,
    )
    section_label(slide, "Week 4 LangGraph + MCP")


def week4_project(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Week 4 项目: AI Workflow Assistant", "用户输入复杂任务，系统规划步骤、调用工具、汇总结果。", ORANGE)
    card(slide, 0.75, 1.35, 3.8, 5.2, "项目结构", [
        "week4-project/workflow.py",
        "week4-project/tools.py",
        "week4-project/models.py",
        "week4-project/main.py",
        "week4-project/templates/index.html",
        "week4-project/eval_cases.json",
        "week4-project/eval_runner.py",
    ], ORANGE)
    card(slide, 4.75, 1.35, 3.8, 5.2, "工作流", [
        "plan_node: LLM 生成计划",
        "execute_node: 调用工具",
        "should_continue: 判断是否继续",
        "summarize_node: 汇总结果",
        "StateGraph: plan → execute loop → summarize",
        "TaskStatus: 保存任务状态",
    ], BLUE)
    card(slide, 8.75, 1.35, 3.8, 5.2, "现在要改", [
        "README 启动命令和导入方式要修",
        "SSE 前端未真正消费 stream",
        "calculator 不应使用 eval",
        "file_reader 要限制目录",
        "web_search 要标明是 mock",
        "补依赖文件和 .env.example",
    ], RED)
    section_label(slide, "Week 4 project")


def eval_system(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "横向能力: 评估体系", "评估不是 Week 4 才做，而是从 Week 2 开始贯穿每个 AI demo。")
    card(slide, 0.75, 1.35, 3.8, 5.2, "评估对象", [
        "Prompt: 输出是否符合格式",
        "结构化抽取: 字段是否完整",
        "RAG 检索: 来源是否命中",
        "RAG 生成: 关键词是否包含",
        "Agent: 工具是否选对",
        "Workflow: 步骤是否合理",
    ], GREEN)
    card(slide, 4.75, 1.35, 3.8, 5.2, "指标", [
        "通过率 pass_rate",
        "工具命中率 tool_rate",
        "关键词命中率 keyword_rate",
        "平均延迟 avg_latency",
        "失败样本数量",
        "token/cost 消耗",
    ], BLUE)
    card(slide, 8.75, 1.35, 3.8, 5.2, "调优闭环", [
        "1. 先跑 baseline",
        "2. 分析失败类型",
        "3. 改 prompt / 工具描述 / 检索参数",
        "4. 重新跑同一评估集",
        "5. 保存 eval_result",
        "6. 写调优记录",
    ], PURPLE)
    section_label(slide, "Evaluation system")


def repo_map(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "仓库结构地图", "你当前仓库已经覆盖大模型应用主链路，下一步是整理成作品集。")
    card(slide, 0.75, 1.35, 2.9, 5.2, "课程目录", [
        "week1-python-ai-basics",
        "week2-prompt-and-langchain",
        "week3-rag-and-fastapi",
        "week4-agent-and-deploy",
        "",
        "作用: 讲义、练习代码、周总结 PPT。",
    ], BLUE)
    card(slide, 3.95, 1.35, 2.9, 5.2, "练习目录", [
        "week1-practice",
        "week2-practice",
        "week3-practice",
        "week4-practice",
        "",
        "作用: 每天小练习，验证单个知识点。",
    ], PURPLE)
    card(slide, 7.15, 1.35, 2.65, 5.2, "项目目录", [
        "week2-project",
        "week3-project",
        "week4-project",
        "",
        "作用: 每周综合作品，可整理为作品集。",
    ], GREEN)
    card(slide, 10.05, 1.35, 2.55, 5.2, "根目录", [
        "learning-plan-v2.md",
        "utils.py",
        "本 PPT",
        "生成脚本",
        "",
        "作用: 路线、公共工具、总结材料。",
    ], ORANGE)
    section_label(slide, "Repo map")


def portfolio(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "作品集整理方式", "最后不要展示 20 个 demo，重点打磨 2-3 个能讲清楚的项目。")
    card(slide, 0.75, 1.35, 3.75, 5.05, "项目 1: AI 数据助手", [
        "展示结构化输出、Tool Use、Pydantic",
        "适合说明: Prompt 工程 + 校验 + 评估",
        "要补: README、示例输入输出、失败案例",
        "关键词: structured output, function calling, eval",
    ], PURPLE)
    card(slide, 4.85, 1.35, 3.75, 5.05, "项目 2: Docs Copilot", [
        "展示 RAG Pipeline、FastAPI、引用来源",
        "适合说明: 检索增强和全栈 AI 应用",
        "要补: 持久化向量库、上传文档、拒答策略",
        "关键词: RAG, embedding, retrieval, SSE",
    ], GREEN)
    card(slide, 8.95, 1.35, 3.75, 5.05, "项目 3: Workflow Assistant", [
        "展示 LangGraph、Agent、工具执行",
        "适合说明: 多步骤工作流和 Agent 边界",
        "要补: 真流式进度、工具安全、部署",
        "关键词: Agent, LangGraph, workflow",
    ], ORANGE)
    section_label(slide, "Portfolio")


def engineering_gaps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "当前项目需要补强的工程问题", "路线是对的，但仓库从学习材料变成作品集还差一层工程整理。")
    add_bullets(slide, 0.9, 1.35, 11.7, 5.35, [
        "1. 统一依赖: 根目录缺少 requirements.txt / pyproject.toml，换环境容易跑不起来。",
        "2. 启动命令: week4-project 的包名和导入方式不够规范，README 命令要修。",
        "3. 真流式: week4-project 声称 SSE，但前端主要还是同步 submit，需要真正消费 /api/stream。",
        "4. 工具安全: calculator 不建议 eval，file_reader 应限制到工作目录或 docs 目录。",
        "5. Mock 声明: web_search 是模拟数据，README 和 UI 里要标清不是实时搜索。",
        "6. README 标准化: 每个项目都要写问题、架构、运行步骤、评估结果、已知问题。",
        "7. 调试输出清理: week2-project/main.py 里有 promptprompt 调试输出，应清理。",
        "8. 评估报告: eval_runner 应保存结果，并对失败样本分类。",
    ], size=14.5, color=BLACK, gap=8)
    section_label(slide, "Engineering gaps")


def next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "下一阶段行动计划", "先把已有项目打磨好，再扩展更多框架或高级能力。")
    card(slide, 0.75, 1.35, 3.8, 5.2, "第 1 阶段\n可运行", [
        "补 requirements.txt",
        "补 .env.example",
        "修 README 启动命令",
        "统一使用 .venv",
        "每个 project 一键运行",
        "删除无关调试输出",
    ], BLUE)
    card(slide, 4.75, 1.35, 3.8, 5.2, "第 2 阶段\n可评估", [
        "扩充 eval_cases",
        "保存 eval_result.json",
        "记录失败原因",
        "统计延迟和 token",
        "对比 prompt 版本",
        "把调优过程写进 README",
    ], PURPLE)
    card(slide, 8.75, 1.35, 3.8, 5.2, "第 3 阶段\n可展示", [
        "整理 2-3 个作品集项目",
        "加架构图和截图",
        "部署至少一个在线 demo",
        "准备面试讲解稿",
        "补前端 AI SDK 加分项",
        "写清下一步优化方向",
    ], GREEN)
    section_label(slide, "Next steps")


def final_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLUE)
    textbox(slide, 0.8, 1.35, 11.8, 0.8, "总结", size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 1.2, 2.35, 11.0, 1.0,
            "这条学习路线是对的。下一步不是再堆更多框架，而是把已有项目打磨到能运行、能评估、能展示。",
            size=23, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, 2.0, 4.15, 9.7, 1.8, [
        "主线: API → Prompt → Structured Output → Eval → RAG → FastAPI → Agent → Deploy",
        "目标: 形成 2-3 个高质量作品集项目",
        "标准: 能讲清架构、边界、失败案例、评估指标和优化过程",
    ], size=16.5, color=PALE_BLUE, gap=12)


def generate():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover(prs)
    how_to_use(prs)
    agenda(prs)
    route_overview(prs)
    route_principles(prs)

    week1_overview(prs)
    week1_details(prs)
    week1_project(prs)

    week2_overview(prs)
    week2_prompt(prs)
    week2_eval_langchain(prs)
    week2_project(prs)

    week3_overview(prs)
    week3_rag(prs)
    week3_fastapi_frontend(prs)
    week3_project(prs)

    week4_overview(prs)
    week4_agent(prs)
    week4_langgraph_mcp(prs)
    week4_project(prs)

    eval_system(prs)
    repo_map(prs)
    portfolio(prs)
    engineering_gaps(prs)
    next_steps(prs)
    final_slide(prs)

    out = Path(__file__).parent / "大模型应用开发学习项目总结.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    generate()
