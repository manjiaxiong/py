"""
第四周学习 PPT 生成脚本
内容：Agent 深入 + LangGraph + MCP + 部署
包含：学习内容 + 面试考点 + 考点答案
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
from pathlib import Path
import io

# ========== 颜色常量 ==========
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x16, 0x65, 0x34)
BG_LIGHT = RGBColor(0xF7, 0xF7, 0xF8)
QA_BLUE = RGBColor(0x1D, 0x4E, 0xD8)
QA_GREEN = RGBColor(0x15, 0x80, 0x3D)


def set_slide_bg(slide, color=BG_LIGHT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, size=13, color=BLACK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = spacing
    return tf


def add_qa_list(slide, left, top, width, height, qa_pairs, q_size=13, a_size=12):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for q, a in qa_pairs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = q
        p.font.size = Pt(q_size)
        p.font.bold = True
        p.font.color.rgb = QA_BLUE
        p.space_before = Pt(10)
        p.space_after = Pt(2)
        p2 = tf.add_paragraph()
        p2.text = a
        p2.font.size = Pt(a_size)
        p2.font.color.rgb = QA_GREEN
        p2.space_after = Pt(6)
    return tf


def inches(val):
    return int(Inches(val))


# ========== 生成 PPT ==========
def generate_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W = inches(13.333)
    H = inches(7.5)
    MARGIN = inches(0.8)
    CONTENT_W = W - 2 * MARGIN

    blank_layout = prs.slide_layouts[6]

    # ========== 封面 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, ACCENT)
    add_text(slide, MARGIN, inches(2.0), CONTENT_W, inches(1.2),
             "第四周学习总结", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN, inches(3.2), CONTENT_W, inches(0.8),
             "Agent 深入 + LangGraph + MCP + 部署", size=24, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN, inches(4.2), CONTENT_W, inches(0.6),
             "Agent Loop | ReAct | LangGraph 工作流 | MCP 协议 | 项目评估调优 | 部署上线",
             size=16, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

    # ========== 目录 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "本周目录", size=28, bold=True, color=ACCENT)
    days = [
        "Day 1: Agent 深入 — Agent Loop、ReAct 模式、停止条件",
        "Day 2: LangGraph — State/Node/Edge、条件边、人工介入",
        "Day 3: MCP — 模型上下文协议、Server/Client 架构",
        "Day 4: 评估与优化 — 评估集、失败分析、调优对比、延迟追踪",
        "Day 5: 部署上线 — Railway/Docker、环境变量、日志中间件",
        "毕业项目: AI Workflow Assistant — LangGraph 多步骤工作流",
        "面试考点: 10 道高频面试题 + 答案",
    ]
    add_bullet_list(slide, MARGIN, inches(1.5), CONTENT_W, inches(5.0),
                    [f"{i+1}. {d}" for i, d in enumerate(days)], size=16, spacing=Pt(12))

    # ========== Day 1: Agent 深入 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "Day 1: Agent 深入", size=28, bold=True, color=ACCENT)
    add_text(slide, MARGIN, inches(1.3), CONTENT_W, inches(0.5),
             "Agent = 让 AI 自主决策调用工具的循环（plan → act → observe → repeat）", size=16, bold=True, color=BLACK)
    items = [
        "Agent vs Chain: Chain = 固定流程(A→B→C), Agent = 动态决策（每步都不同）",
        "ReAct 模式: Thought(推理分析) → Action(选择工具) → Observation(工具返回) → 循环 → Final Answer",
        "ReAct Prompt 模板: 明确定义 Thought/Action/Action Input/Final Answer 格式",
        "工具注册表: dict 映射工具名到函数，就像路由表",
        "停止条件（必须有）: Final Answer / max_iterations / 工具失败重试耗尽",
        "不是所有任务都适合 Agent — 简单任务直接调用更快更稳定更省钱",
        "",
        "适用 Agent: 多步骤研究、不确定路径、需要多个工具",
        "不适合 Agent: 简单翻译、固定格式输出、对延迟敏感",
    ]
    add_bullet_list(slide, MARGIN, inches(2.0), CONTENT_W, inches(5.0), items, size=14, spacing=Pt(6))

    # ========== Day 1: SimpleAgent 类 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "Day 1: 手写 SimpleAgent", size=28, bold=True, color=ACCENT)
    half_w = CONTENT_W // 2 - inches(0.3)

    add_text(slide, MARGIN, inches(1.3), half_w, inches(0.5),
             "类结构", size=18, bold=True, color=ACCENT2)
    items_left = [
        "__init__: client, model, tools, max_iterations, max_retries",
        "run(task): Agent 循环入口",
        "  for i in range(max_iterations):",
        "    1. 构造 ReAct Prompt",
        "    2. 调 LLM 获取下一步",
        "    3. parse_react_response() 解析",
        "    4. 如果 Final Answer → 返回",
        "    5. 如果 Action → _execute_tool()",
        "    6. 更新 history 继续循环",
        "_execute_tool: 调用工具 + 重试",
        "_parse_response: 正则提取 Action/Input",
    ]
    add_bullet_list(slide, MARGIN, inches(2.0), half_w, inches(4.5), items_left, size=12, spacing=Pt(4))

    right_x = MARGIN + half_w + inches(0.6)
    add_text(slide, right_x, inches(1.3), half_w, inches(0.5),
             "JS 类比", size=18, bold=True, color=ACCENT2)
    items_right = [
        "class Agent = class Agent { constructor(), async run() }",
        "while i < max = while (attempts < MAX)",
        "TOOLS dict = const tools = { calc: fn, search: fn }",
        "re.search = text.match(/pattern/)",
        "try/except + retry = try/catch + for loop",
        "history += observation = context.push(observation)",
        "parse_react_response = parseResponse(text)",
    ]
    add_bullet_list(slide, right_x, inches(2.0), half_w, inches(4.5), items_right, size=12, spacing=Pt(4))

    # ========== Day 2: LangGraph 核心概念 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "Day 2: LangGraph 核心概念", size=28, bold=True, color=ACCENT)
    add_text(slide, MARGIN, inches(1.3), CONTENT_W, inches(0.4),
             "为什么需要 LangGraph: LangChain Agent 黑盒不可控, LangGraph 图结构清晰可见", size=16, bold=True, color=BLACK)

    half_w = CONTENT_W // 2 - inches(0.3)
    add_text(slide, MARGIN, inches(2.0), half_w, inches(0.5),
             "四大核心概念", size=18, bold=True, color=ACCENT2)
    items_left = [
        "State(TypedDict): 共享数据对象，所有节点都能读写",
        "  JS 类比: React useState / Redux store",
        "",
        "Node(函数): 接收 state，返回更新字典",
        "  JS 类比: reducer (state) => newState",
        "",
        "Edge(边): A → B 固定连接，无条件跳转",
        "  JS 类比: 直接函数调用 return next()",
        "",
        "Conditional Edge: if/else 条件路由",
        "  JS 类比: switch/case 路由",
    ]
    add_bullet_list(slide, MARGIN, inches(2.8), half_w, inches(4.0), items_left, size=13, spacing=Pt(3))

    right_x = MARGIN + half_w + inches(0.6)
    add_text(slide, right_x, inches(2.0), half_w, inches(0.5),
             "代码骨架", size=18, bold=True, color=ACCENT2)
    items_right = [
        "graph = StateGraph(WorkflowState)",
        "graph.add_node('research', research_node)",
        "graph.add_node('outline', outline_node)",
        "graph.add_edge('research', 'outline')",
        "",
        "# 条件边",
        "def should_continue(state):",
        "    if state['score'] >= 7:",
        "        return 'end'",
        "    return 'retry'",
        "graph.add_conditional_edges(",
        "    'review', should_continue,",
        "    {'end': END, 'retry': 'draft'},",
        ")",
        "app = graph.compile()",
        "result = app.invoke(initial_state)",
    ]
    add_bullet_list(slide, right_x, inches(2.8), half_w, inches(4.0), items_right, size=11, spacing=Pt(2))

    # ========== Day 2: 人工介入 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "Day 2: 人工介入 (Human-in-the-loop)", size=28, bold=True, color=ACCENT)
    items = [
        "在关键步骤暂停，等人工审核后继续 — 生产环境必备能力",
        "",
        "实现方式:",
        "  1. MemorySaver() — 内存中保存状态快照（检查点）",
        "  2. compile(interrupt_before=['review']) — 在 review 前暂停",
        "  3. 第一次 invoke() — 运行到 interrupt 点自动停止",
        "  4. Command(resume='approved') — 恢复并传入人工反馈",
        "",
        "JS 类比:",
        "  const result = await workflow.runUntil('review');",
        "  const feedback = await getUserInput();",
        "  await workflow.resume(feedback);",
        "",
        "使用场景:",
        "  - 内容审核（AI 生成 → 人工确认 → 发布）",
        "  - 敏感操作审批（AI 建议 → 人工决定）",
        "  - 质量保证（草稿 → 人工修改 → 最终稿）",
    ]
    add_bullet_list(slide, MARGIN, inches(1.5), CONTENT_W, inches(5.5), items, size=14, spacing=Pt(6))

    # ========== Day 3: MCP ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "Day 3: MCP 模型上下文协议", size=28, bold=True, color=ACCENT)
    add_text(slide, MARGIN, inches(1.3), CONTENT_W, inches(0.4),
             "MCP = AI 工具的 USB 接口，统一协议，即插即用", size=16, bold=True, color=BLACK)

    half_w = CONTENT_W // 2 - inches(0.3)
    add_text(slide, MARGIN, inches(2.0), half_w, inches(0.5),
             "MCP vs Function Calling", size=18, bold=True, color=ACCENT2)
    items_left = [
        "Function Calling: 工具硬编码在 prompt/代码里",
        "  改工具 → 改代码 → 重新部署",
        "MCP: 工具由 Server 动态提供",
        "  改工具 → 只改 Server → Client 自动发现",
        "",
        "架构: Client(AI App) ←→ Server(工具提供方)",
        "",
        "协议流程:",
        "  1. 连接 Server",
        "  2. tools/list → 获取工具列表",
        "  3. tools/call → 调用指定工具",
    ]
    add_bullet_list(slide, MARGIN, inches(2.8), half_w, inches(4.0), items_left, size=13, spacing=Pt(6))

    right_x = MARGIN + half_w + inches(0.6)
    add_text(slide, right_x, inches(2.0), half_w, inches(0.5),
             "JS 类比", size=18, bold=True, color=ACCENT2)
    items_right = [
        "MCP 协议 = REST API 标准",
        "MCP Server = Express 服务器（提供工具）",
        "MCP Client = fetch/axios（消费工具）",
        "tools/list = GET /api/tools（发现）",
        "tools/call = POST /api/tools/:name（调用）",
        "",
        "以前: 每个组件手动 fetch 不同 API",
        "现在: 浏览器 Extension API 写一次到处用",
        "",
        "工具定义 = OpenAPI Schema",
        "stdio 通信 = 进程间通信(IPC)",
    ]
    add_bullet_list(slide, right_x, inches(2.8), half_w, inches(4.0), items_right, size=13, spacing=Pt(6))

    # ========== Day 4: 评估与优化 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "Day 4: 评估与优化", size=28, bold=True, color=ACCENT)
    items = [
        "评估集: 固定输入 + 期望关键词 → 跑通过率和延迟",
        "",
        "失败原因分类:",
        "  Prompt 不清晰 → 回答跑题 → 改 system prompt",
        "  工具描述歧义 → Agent 选错 → 改 tool description",
        "  检索不相关 → RAG 答非所问 → 调 chunk_size/top_k",
        "  模型能力不足 → 简单题出错 → 换更强模型",
        "",
        "多轮调优: baseline → 加 system prompt → 对比通过率 → 选最优",
        "",
        "延迟追踪: @track_latency 装饰器记录每次调用耗时",
        "优化策略: 减少 max_tokens | 缓存重复查询 | 并行调用 | 用更快模型",
        "",
        "调优报告: Markdown 格式，对比表 + 优化建议",
    ]
    add_bullet_list(slide, MARGIN, inches(1.5), CONTENT_W, inches(5.5), items, size=14, spacing=Pt(6))

    # ========== Day 5: 部署 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "Day 5: 部署上线", size=28, bold=True, color=ACCENT)

    half_w = CONTENT_W // 2 - inches(0.3)
    add_text(slide, MARGIN, inches(1.5), half_w, inches(0.5),
             "部署方案", size=18, bold=True, color=ACCENT2)
    items_left = [
        "Railway: 最简单，push 就部署 ($5/月起)",
        "  JS 类比: Vercel for Python",
        "",
        "Docker: 打包环境，一致性强",
        "  Dockerfile + docker-compose.yml",
        "  FROM python:3.11-slim → pip install → uvicorn",
        "",
        "部署文件清单:",
        "  requirements.txt — 依赖列表",
        "  Procfile — 启动命令",
        "  Dockerfile — Docker 镜像",
        "  .env.example — 环境变量模板",
    ]
    add_bullet_list(slide, MARGIN, inches(2.2), half_w, inches(4.5), items_left, size=13, spacing=Pt(5))

    right_x = MARGIN + half_w + inches(0.6)
    add_text(slide, right_x, inches(1.5), half_w, inches(0.5),
             "安全与监控", size=18, bold=True, color=ACCENT2)
    items_right = [
        "环境变量安全:",
        "  .env 加入 .gitignore（绝不提交）",
        "  提供 .env.example（隐藏值）",
        "  线上用平台环境变量设置",
        "",
        "日志中间件:",
        "  @app.middleware('http')",
        "  记录: method, path, status, latency_ms",
        "  JS 类比: Express logging middleware",
        "",
        "日志级别: DEBUG → INFO → WARNING → ERROR",
        "  JS 类比: console.debug/log/warn/error",
    ]
    add_bullet_list(slide, right_x, inches(2.2), half_w, inches(4.5), items_right, size=13, spacing=Pt(5))

    # ========== 毕业项目: AI Workflow Assistant ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "毕业项目: AI Workflow Assistant", size=28, bold=True, color=ACCENT)
    add_text(slide, MARGIN, inches(1.3), CONTENT_W, inches(0.4),
             "用户输入任务描述，LangGraph 自动规划并执行多步骤工作流", size=16, bold=True, color=BLACK)
    items = [
        "workflow.py — LangGraph StateGraph（plan → execute loop → summarize → END）",
        "tools.py — 5 个工具（web_search, calculator, file_reader, text_summarize, get_current_time）",
        "models.py — Pydantic 模型（TaskSubmitRequest, StepInfo, TaskStatusResponse）",
        "main.py — FastAPI + SSE 流式进度 + 统一响应格式 {code, msg, data}",
        "templates/index.html — 步骤卡片 UI（运行中/成功/失败状态）",
        "eval_cases.json — 8 条评估用例 + eval_runner.py 评估脚本",
        "",
        "工作流流程:",
        "  plan_node(LLM 生成计划) → execute_node(调用工具) → 条件判断(继续/结束) → summarize_node(LLM 总结)",
        "",
        "运行: uvicorn week4-project.main:app --reload → 访问 localhost:8000",
    ]
    add_bullet_list(slide, MARGIN, inches(2.0), CONTENT_W, inches(5.0), items, size=13, spacing=Pt(5))

    # ========== 核心架构: LangGraph 工作流 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "核心架构: LangGraph 工作流", size=28, bold=True, color=ACCENT)

    fig, ax = plt.subplots(1, 1, figsize=(11, 4.0))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 4)
    ax.axis("off")
    fig.patch.set_facecolor("#F7F7F8")

    font_prop = {"fontfamily": "Microsoft YaHei", "fontsize": 10}

    boxes = [
        (1.0, 2.5, "用户任务\n描述", "#DBEAFE"),
        (3.0, 2.5, "plan_node\nLLM 规划", "#E0E7FF"),
        (5.5, 2.5, "execute_node\n调用工具", "#EDE9FE"),
        (8.0, 2.5, "条件判断\n继续/结束?", "#FEF3C7"),
        (10.0, 2.5, "summarize\nLLM 总结", "#DCFCE7"),
    ]

    for x, y, text, color in boxes:
        hex_color = color.lstrip("#")
        rgb = tuple(int(hex_color[i:i+2], 16) / 255.0 for i in (0, 2, 4))
        rect = FancyBboxPatch((x, y), 1.6, 1.0, boxstyle="round,pad=0.1",
                              facecolor=rgb, edgecolor="#94A3B8", linewidth=1)
        ax.add_patch(rect)
        ax.text(x + 0.8, y + 0.5, text, ha="center", va="center", **font_prop)

    # 箭头
    for x1, x2 in [(2.0, 3.0), (4.2, 5.5), (6.7, 8.0), (9.2, 10.0)]:
        ax.annotate("", xy=(x2, 3.0), xytext=(x1 + 0.6, 3.0),
                    arrowprops=dict(arrowstyle="->", color="#64748B", lw=1.5))

    # 循环箭头（从条件判断回到执行）
    ax.annotate("", xy=(5.5, 2.2), xytext=(8.8, 2.2),
                arrowprops=dict(arrowstyle="->", color="#7C3AED", lw=1.5, linestyle="--"))
    ax.text(7.15, 1.9, "继续执行", ha="center", fontsize=9,
            fontfamily="Microsoft YaHei", color="#7C3AED")

    # 工具标注
    ax.text(6.3, 2.0, "tools.py", ha="center", fontsize=8,
            fontfamily="Microsoft YaHei", color="#2563EB", fontweight="bold")

    # 标注
    ax.text(5.5, 3.8, "LangGraph StateGraph 工作流", ha="center", fontsize=12,
            fontfamily="Microsoft YaHei", fontweight="bold", color="#2563EB")
    ax.text(5.5, 1.5, "虚线 = 条件边（根据判断决定走向）", ha="center", fontsize=9,
            fontfamily="Microsoft YaHei", color="#6B7280")

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="#F7F7F8")
    plt.close(fig)
    buf.seek(0)
    slide.shapes.add_picture(buf, Emu(MARGIN), Emu(inches(1.2)), Emu(CONTENT_W), Emu(inches(5.0)))

    # ========== 面试考点 1-5 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.4), CONTENT_W, inches(0.7),
             "面试考点 (1/2)", size=28, bold=True, color=ACCENT)
    qa_pairs_1 = [
        ("Q1: Agent 和 Chain 有什么区别？什么场景适合用 Agent？",
         "A: Chain = 固定流程（A→B→C），步骤写死不可变。Agent = 动态决策，"
         "每步由 LLM 决定下一步做什么。适合 Agent 的场景：多步骤研究任务、"
         "不确定需要哪些步骤、需要调用多个外部工具。不适合的场景：简单翻译/摘要"
         "（一次调用就够）、固定格式输出（Chain 更稳定）、对延迟敏感的场景。"),

        ("Q2: 请解释 ReAct 模式的工作原理",
         "A: ReAct = Reasoning + Acting。流程：Thought（分析当前情况）→ Action"
         "（选择工具）→ Observation（查看工具返回结果）→ 回到 Thought 继续推理。"
         "当 LLM 判断已有足够信息时输出 Final Answer。关键是用明确的 Prompt 格式"
         "让 LLM 按固定格式输出，然后用正则解析并执行对应工具。"),

        ("Q3: Agent 为什么必须有停止条件？有哪些停止方式？",
         "A: Agent 本质是 while 循环，没有停止条件会无限循环（LLM 可能重复选择同一工具）。"
         "停止方式：1) LLM 输出 Final Answer（正常结束）；2) 达到 max_iterations（强制结束）；"
         "3) 工具调用失败且重试耗尽（报错结束）。max_iterations 通常设为 5-10。"),

        ("Q4: LangGraph 和 LangChain Agent 的主要区别是什么？",
         "A: LangChain Agent 是黑盒，流程由 LLM 自己决定，不可控、无法暂停、难以调试。"
         "LangGraph 是图结构，每个 Node 做什么、Edge 怎么连接都显式定义。"
         "LangGraph 支持人工介入（interrupt_before），状态透明（State 对象可观察），"
         "更适合生产环境的复杂工作流。"),

        ("Q5: LangGraph 中的 State、Node、Edge 分别是什么？JS 类比是什么？",
         "A: State = 共享数据对象（TypedDict），JS 类比 React useState/Redux store。"
         "Node = 执行单元函数，接收 state 返回更新，JS 类比 reducer。"
         "Edge = 节点间的连接，A→B 无条件跳转，JS 类比直接函数调用。"
         "Conditional Edge = 条件路由，根据 state 决定下一个节点，JS 类比 switch/case。"),
    ]
    add_qa_list(slide, MARGIN, inches(1.2), CONTENT_W, inches(5.8), qa_pairs_1)

    # ========== 面试考点 6-10 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.4), CONTENT_W, inches(0.7),
             "面试考点 (2/2)", size=28, bold=True, color=ACCENT)
    qa_pairs_2 = [
        ("Q6: MCP 是什么？和 Function Calling 有什么区别？",
         "A: MCP（Model Context Protocol）是 Anthropic 提出的开放协议，"
         "让 AI 应用能即插即用地接入外部工具。区别：Function Calling 的工具"
         "定义硬编码在 prompt/代码里，改工具需要改代码重新部署。MCP 的工具由 Server"
         "动态提供，Client 通过 tools/list 自动发现，改工具只需改 Server，"
         "Client 不用动。类比：Function Calling = 手动 import，MCP = npm install。"),

        ("Q7: 怎么评估一个 AI Agent 系统的效果？",
         "A: 设计评估集（固定输入 + 期望输出关键词），运行后检查：1) 工具命中率 — "
         "期望的工具是否被调用；2) 关键词通过率 — 结果是否包含期望关键词；3) 延迟和成本"
         "— 每次调用耗时和 token 消耗。改 prompt/工具描述/参数前后都跑评估，对比数据选最优。"),

        ("Q8: 部署 Python 项目有哪些常见方案？各自适用场景？",
         "A: Railway: 最简单，连接 GitHub 仓库自动部署，$5/月起，适合快速上线 demo。"
         "Render: 有免费额度，类似 Railway。Docker: 打包整个运行环境，一致性强，"
         "适合生产环境。部署前需要准备: requirements.txt、Procfile、Dockerfile、.env.example。"),

        ("Q9: 线上项目为什么需要日志中间件？怎么实现？",
         "A: 日志是线上出问题时唯一的线索。中间件在每次请求前后拦截，记录方法、路径、"
         "状态码、耗时，方便定位慢请求和错误。FastAPI 实现: "
         "@app.middleware('http') 装饰器，time.time() 记录耗时，"
         "logging.info(json.dumps(...)) 输出结构化日志。"),

        ("Q10: AI 系统调优的一般流程是什么？",
         "A: 1) 设计评估集（覆盖核心场景 + 边界情况）；2) 跑基线评估记录通过率和延迟；"
         "3) 分析失败用例，分类原因（prompt/工具/检索/模型）；4) 多轮调优对比"
         "（改 prompt / 换参数 / 加 system prompt），同一评估集对比数据；"
         "5) 选最优配置，生成调优报告；6) 每次新增功能或改 prompt 都跑回归测试。"),
    ]
    add_qa_list(slide, MARGIN, inches(1.2), CONTENT_W, inches(5.8), qa_pairs_2)

    # ========== 总结页 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, ACCENT)
    add_text(slide, MARGIN, inches(2.0), CONTENT_W, inches(1.0),
             "Week 4 Complete!", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN, inches(3.2), CONTENT_W, inches(0.6),
             "从 Agent 到部署上线，4 周学习圆满收官", size=20, color=WHITE, align=PP_ALIGN.CENTER)
    summary_items = [
        "Agent: Agent Loop + ReAct + 停止条件 + 错误回退",
        "LangGraph: State/Node/Edge + 条件边 + 人工介入",
        "MCP: 统一工具协议 + Server/Client + 自动发现",
        "评估: 评估集 + 失败分析 + 调优对比 + 延迟追踪",
        "部署: Railway/Docker + 环境变量安全 + 日志中间件",
        "毕业项目: AI Workflow Assistant (LangGraph + 5 工具 + FastAPI)",
    ]
    add_bullet_list(slide, inches(2.5), inches(4.2), inches(8.333), inches(3.0),
                    summary_items, size=16, color=RGBColor(0xBF, 0xDB, 0xFE), spacing=Pt(10))

    # ========== 保存 ==========
    output = Path(__file__).parent / "week4_summary.pptx"
    prs.save(output)
    print(f"PPT 已保存到: {output}")
    return output


if __name__ == "__main__":
    generate_ppt()
