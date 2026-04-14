"""
第三周学习 PPT 生成脚本
内容：RAG + FastAPI + 全栈整合
包含：学习内容 + 面试考点 + 考点答案
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import numpy as np
from pathlib import Path
import io

# ========== 颜色常量 ==========
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x16, 0x65, 0x34)
BG_LIGHT = RGBColor(0xF7, 0xF7, 0xF8)
QA_BLUE = RGBColor(0x1D, 0x4E, 0xD8)
QA_GREEN = RGBColor(0x15, 0x80, 0x3D)


def set_slide_bg(slide, color=BG_LIGHT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, size=13, color=BLACK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = spacing
    return tf


def add_qa_list(slide, left, top, width, height, qa_pairs, q_size=13, a_size=12):
    """添加 Q&A 列表，Q 蓝色加粗，A 绿色"""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for q, a in qa_pairs:
        # Q
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = q
        p.font.size = Pt(q_size)
        p.font.bold = True
        p.font.color.rgb = QA_BLUE
        p.space_before = Pt(10)
        p.space_after = Pt(2)
        # A
        p2 = tf.add_paragraph()
        p2.text = a
        p2.font.size = Pt(a_size)
        p2.font.color.rgb = QA_GREEN
        p2.space_after = Pt(6)
    return tf


# ========== 单位辅助 ==========
def inches(val):
    return int(Inches(val))


# ========== 生成 PPT ==========
def generate_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W = inches(13.333)
    H = inches(7.5)
    MARGIN = inches(0.8)
    CONTENT_W = W - 2 * MARGIN

    blank_layout = prs.slide_layouts[6]

    # ========== 封面 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, ACCENT)
    add_text(slide, MARGIN, inches(2.0), CONTENT_W, inches(1.2),
             "第三周学习总结", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN, inches(3.2), CONTENT_W, inches(0.8),
             "RAG + FastAPI + 全栈整合", size=24, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN, inches(4.2), CONTENT_W, inches(0.6),
             "检索增强生成 | Python Web 框架 | SSE 流式输出 | 前端对接",
             size=16, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

    # ========== 目录 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "本周目录", size=28, bold=True, color=ACCENT)
    days = [
        "Day 1: RAG 基础概念 — 分块、Embedding、向量检索",
        "Day 2: 最小 RAG Pipeline — 加载 -> 分块 -> 存储 -> 检索 -> 生成",
        "Day 3: RAG 优化 — chunk_size、top_k、阈值、Rerank、评估",
        "Day 4: FastAPI 后端 — 路由、Pydantic、SSE、CORS",
        "Day 5: 前端对接 — SSE 流式接收、聊天界面",
        "周项目: Docs Copilot — AI 文档问答助手",
        "面试考点: 10 道高频面试题 + 答案",
    ]
    add_bullet_list(slide, MARGIN, inches(1.5), CONTENT_W, inches(5.0),
                    [f"{i+1}. {d}" for i, d in enumerate(days)], size=16, spacing=Pt(12))

    # ========== Day 1: RAG 基础概念 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "Day 1: RAG 基础概念", size=28, bold=True, color=ACCENT)
    add_text(slide, MARGIN, inches(1.3), CONTENT_W, inches(0.5),
             "RAG = Retrieval-Augmented Generation（检索增强生成）", size=16, bold=True, color=BLACK)
    items = [
        "核心思想: AI 不靠记忆回答，先从文档中搜索相关内容，再基于搜索结果回答",
        "类比: 闭卷考试(直接问AI) vs 开卷考试(RAG) vs 提前背书(微调)",
        "RAG 优势: 便宜、实时更新、可控、能处理私有数据、减少幻觉",
        "完整流程:",
        "  离线阶段: 文档 -> 分块(Chunking) -> 向量化(Embedding) -> 存入向量库",
        "  在线阶段: 用户提问 -> 向量化 -> 向量库检索 -> 取 top-k -> 拼接 prompt -> LLM 生成",
    ]
    add_bullet_list(slide, MARGIN, inches(2.0), CONTENT_W, inches(4.5), items, size=14, spacing=Pt(8))

    # ========== Day 1: Chunking + Embedding ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "Day 1: Chunking + Embedding", size=28, bold=True, color=ACCENT)

    # 左半：Chunking
    half_w = CONTENT_W // 2 - inches(0.3)
    add_text(slide, MARGIN, inches(1.3), half_w, inches(0.5),
             "Chunking（文档分块）", size=18, bold=True, color=ACCENT2)
    chunk_items = [
        "为什么分块: token 限制 + 精准匹配",
        "固定字数: 每 300 字切一块（简单粗暴）",
        "按段落: 用 \\n\\n 分割（保留语义）",
        "递归切分: 先段落 -> 句子 -> 字数（最佳实践）",
        "overlap: 相邻块重叠 50-100 字防切断",
        "推荐: 300-500 字，视场景调整",
    ]
    add_bullet_list(slide, MARGIN, inches(2.0), half_w, inches(4.5), chunk_items, size=13, spacing=Pt(6))

    # 右半：Embedding
    right_x = MARGIN + half_w + inches(0.6)
    add_text(slide, right_x, inches(1.3), half_w, inches(0.5),
             "Embedding（文本向量化）", size=18, bold=True, color=ACCENT2)
    embed_items = [
        "本质: 文本 -> 一组数字（384 维向量）",
        "特性: 语义相近 -> 向量距离近",
        '例: "React 组件" 和 "Vue 组件" 向量很接近',
        "余弦相似度: 1=完全相同, 0=无关",
        "L2 距离: 0=完全相同, 越大越远",
        "Chroma 默认: all-MiniLM-L6-v2（本地免费）",
    ]
    add_bullet_list(slide, right_x, inches(2.0), half_w, inches(4.5), embed_items, size=13, spacing=Pt(6))

    # ========== Day 2: RAG Pipeline ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "Day 2: 最小 RAG Pipeline", size=28, bold=True, color=ACCENT)
    items = [
        "Step 1: 加载文档 — Path(dir).glob('*.md') 读取所有 Markdown 文件",
        "Step 2: 文档分块 — 按段落 split('\\n\\n') + 过滤空段落 + overlap 重叠",
        "Step 3: 存入向量库 — collection.add(documents, ids, metadatas)",
        "Step 4: 检索测试 — collection.query(query_texts=[问题], n_results=3)",
        "Step 5: RAG 问答 — 检索结果拼接 context -> 构造 prompt -> 调 LLM 生成",
        "",
        "ChromaDB 核心 API:",
        "  add() 添加文档 | query() 语义检索 | get() 按 ID 获取",
        "  delete() 删除 | update() 更新 | count() 统计总数",
        "",
        "两种模式:",
        "  Client() — 内存模式，重启丢失（学习用）",
        "  PersistentClient(path='./db') — 持久化到磁盘",
    ]
    add_bullet_list(slide, MARGIN, inches(1.5), CONTENT_W, inches(5.0), items, size=14, spacing=Pt(5))

    # ========== Day 3: RAG 优化 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "Day 3: RAG 优化", size=28, bold=True, color=ACCENT)
    items = [
        "1. chunk_size 实验: 200/500/1000 字符对比，500 通常最佳",
        "2. top_k 调整: k=1 不完整, k=3 推荐, k=5+ 可能有噪音",
        "3. 引用来源: metadata 中存 source 文件名，检索时一起返回",
        "4. 相似度阈值: 距离 > 阈值 -> 拒绝回答（防止基于不相关内容编造）",
        "5. Rerank 概念:",
        "   Bi-encoder(当前): 文档和查询分别编码 -> 快但精度一般",
        "   Cross-encoder(Rerank): 查询+文档一起输入 -> 慢但精度高",
        "   实际: 先 Bi-encoder 召回 20 条 -> Cross-encoder 精排取 3 条",
        "6. 评估 RAG（最重要）:",
        "   检索评估: 期望来源是否在 top-k 中？(命中率)",
        "   生成评估: AI 回答是否包含期望关键词？(通过率)",
        "   分开评估: 检索错 -> 优化分块/embedding; 生成错 -> 优化 prompt",
    ]
    add_bullet_list(slide, MARGIN, inches(1.5), CONTENT_W, inches(5.5), items, size=13, spacing=Pt(5))

    # ========== Day 4: FastAPI 基础 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "Day 4: FastAPI 后端", size=28, bold=True, color=ACCENT)
    add_text(slide, MARGIN, inches(1.3), CONTENT_W, inches(0.4),
             "FastAPI = Python 版 Express.js，自带 API 文档 + 类型校验", size=16, bold=True, color=BLACK)
    items = [
        "创建应用: app = FastAPI() | 等价 const app = express()",
        "GET 路由: @app.get('/path') | 路径参数 {name} | 查询参数 page=1",
        "POST 路由: @app.post('/path') + Pydantic Model = 自动解析+校验",
        "Pydantic: BaseModel + Field(min_length, ge, le) = 内置 Zod",
        "响应模型: response_model 自动过滤字段（防泄露敏感数据）",
        "RAG 接口: POST /index 索引文档 + POST /ask RAG 问答",
        "统一响应: success()/error() 返回 {code, msg, data} 格式",
        "SSE 流式: StreamingResponse + yield（打字机效果）",
        "CORS: app.add_middleware(CORSMiddleware, allow_origins=[...])",
        "启动: uvicorn app:app --reload | 访问 /docs 看自动生成的 Swagger",
    ]
    add_bullet_list(slide, MARGIN, inches(2.0), CONTENT_W, inches(5.0), items, size=13, spacing=Pt(5))

    # ========== Day 4: Express vs FastAPI 对比 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "Day 4: Express vs FastAPI 速查", size=28, bold=True, color=ACCENT)
    comparisons = [
        ("app.get('/path', handler)", "@app.get('/path')"),
        ("req.params.id", "函数参数 (path param)"),
        ("req.query.page", "函数参数 + 默认值 (query param)"),
        ("req.body + Zod", "Pydantic Model 参数"),
        ("res.json({...})", "return {...}"),
        ("res.status(201).json(...)", "JSONResponse(status_code=201, ...)"),
        ("app.use(cors())", "app.add_middleware(CORSMiddleware, ...)"),
        ("app.use(express.static(...))", "app.mount('/static', StaticFiles(...))"),
        ("res.write() / SSE", "StreamingResponse / EventSourceResponse"),
        ("nodemon", "uvicorn --reload"),
        ("手动配 Swagger", "自动生成 /docs"),
    ]
    # 表头
    add_text(slide, MARGIN, inches(1.5), CONTENT_W // 2, inches(0.4),
             "Express.js", size=15, bold=True, color=ACCENT)
    add_text(slide, MARGIN + CONTENT_W // 2, inches(1.5), CONTENT_W // 2, inches(0.4),
             "FastAPI", size=15, bold=True, color=ACCENT2)
    # 表内容
    y_start = inches(2.0)
    row_h = inches(0.38)
    for i, (express, fastapi) in enumerate(comparisons):
        y = y_start + i * row_h
        add_text(slide, MARGIN, y, CONTENT_W // 2, row_h, express, size=12, color=GRAY)
        add_text(slide, MARGIN + CONTENT_W // 2, y, CONTENT_W // 2, row_h, fastapi, size=12, color=BLACK)

    # ========== Day 5: 前端对接 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "Day 5: 前端对接", size=28, bold=True, color=ACCENT)
    items = [
        "方案: 纯 HTML + CSS + JS（学习阶段，零依赖无需构建工具）",
        "",
        "SSE 完整链路:",
        "  后端: async def generator() + yield + EventSourceResponse",
        "  协议: event: token\\ndata: {\"text\": \"xxx\"}\\n\\n",
        "  前端: fetch + ReadableStream + TextDecoder（不用 EventSource，因为需要 POST）",
        "",
        "前端技术点:",
        "  fetch + POST 发送问题",
        "  response.body.getReader() 读取 SSE 流",
        "  TextDecoder 二进制转字符串",
        "  innerHTML += 逐字追加（打字机效果）",
        "  AbortController 取消请求",
        "  CSS Variables 主题管理",
        "",
        "引用来源: 折叠面板展示检索到的原文片段 + 文件名",
    ]
    add_bullet_list(slide, MARGIN, inches(1.5), CONTENT_W, inches(5.5), items, size=13, spacing=Pt(4))

    # ========== 周项目: Docs Copilot ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "周项目: Docs Copilot", size=28, bold=True, color=ACCENT)
    add_text(slide, MARGIN, inches(1.3), CONTENT_W, inches(0.4),
             "AI 文档问答助手 — 整合 RAG + FastAPI + 前端", size=16, bold=True, color=BLACK)
    items = [
        "rag.py — RAGPipeline 类封装（分块 / 索引 / 检索 / 问答 / 阈值过滤）",
        "models.py — Pydantic 数据模型（IndexRequest, AskRequest, AskResponse...）",
        "main.py — FastAPI 服务入口（GET / + POST /api/ask + SSE 流式）",
        "indexer.py — 文档索引脚本（读 docs/ 下所有 .md -> 分块 -> 存入向量库）",
        "templates/index.html — 聊天界面（SSE 流式接收 + 引用来源折叠）",
        "eval_cases.json — 10 条评估用例（含拒答测试）",
        "eval_runner.py — 评估脚本（检索命中率 + 生成通过率 分开评估）",
        "",
        "运行方式:",
        "  python indexer.py  -> 索引文档",
        "  uvicorn main:app --reload  -> 启动服务",
        "  python eval_runner.py  -> 跑评估",
    ]
    add_bullet_list(slide, MARGIN, inches(2.0), CONTENT_W, inches(5.0), items, size=13, spacing=Pt(5))

    # ========== 本周核心架构图 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.5), CONTENT_W, inches(0.7),
             "核心架构: RAG 全栈流程", size=28, bold=True, color=ACCENT)

    fig, ax = plt.subplots(1, 1, figsize=(11, 4.5))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 5)
    ax.axis("off")
    fig.patch.set_facecolor("#F7F7F8")

    font_prop = {"fontfamily": "Microsoft YaHei", "fontsize": 10}

    boxes = [
        (0.5, 3.5, "docs/\n.md 文件", "#DBEAFE"),
        (2.5, 3.5, "Chunking\n分块", "#E0E7FF"),
        (4.5, 3.5, "Embedding\n向量化", "#EDE9FE"),
        (6.5, 3.5, "ChromaDB\n存储", "#FEF3C7"),
        (0.5, 1.0, "用户提问", "#DCFCE7"),
        (2.5, 1.0, "Embedding\n向量化", "#EDE9FE"),
        (4.5, 1.0, "向量检索\ntop-k", "#FEF3C7"),
        (6.5, 1.0, "拼接\nPrompt", "#FFE4E6"),
        (8.5, 1.0, "LLM\n生成回答", "#DBEAFE"),
        (8.5, 3.5, "FastAPI\nSSE 流式", "#DCFCE7"),
    ]

    for x, y, text, color in boxes:
        hex_color = color.lstrip("#")
        rgb = tuple(int(hex_color[i:i+2], 16) / 255.0 for i in (0, 2, 4))
        rect = FancyBboxPatch((x, y), 1.6, 1.0, boxstyle="round,pad=0.1",
                              facecolor=rgb, edgecolor="#94A3B8", linewidth=1)
        ax.add_patch(rect)
        ax.text(x + 0.8, y + 0.5, text, ha="center", va="center", **font_prop)

    # 箭头 - 离线阶段
    for x1, x2 in [(1.5, 2.5), (3.5, 4.5), (5.5, 6.5), (7.5, 8.5)]:
        ax.annotate("", xy=(x2, 4.0), xytext=(x1 + 0.6, 4.0),
                    arrowprops=dict(arrowstyle="->", color="#64748B", lw=1.5))
    # 箭头 - 在线阶段
    for x1, x2 in [(1.5, 2.5), (3.5, 4.5), (5.5, 6.5), (7.5, 8.5)]:
        ax.annotate("", xy=(x2, 1.5), xytext=(x1 + 0.6, 1.5),
                    arrowprops=dict(arrowstyle="->", color="#64748B", lw=1.5))
    # 箭头 - 检索连到 ChromaDB
    ax.annotate("", xy=(6.5, 3.5), xytext=(5.3, 2.0),
                arrowprops=dict(arrowstyle="->", color="#2563EB", lw=1.5, linestyle="--"))
    # 箭头 - FastAPI 到前端
    ax.annotate("", xy=(10.1, 3.5), xytext=(10.1, 2.0),
                arrowprops=dict(arrowstyle="->", color="#2563EB", lw=1.5, linestyle="--"))

    # 标注
    ax.text(5.5, 4.8, "离线阶段（建索引）", ha="center", fontsize=11,
            fontfamily="Microsoft YaHei", fontweight="bold", color="#2563EB")
    ax.text(5.5, 0.5, "在线阶段（回答问题）", ha="center", fontsize=11,
            fontfamily="Microsoft YaHei", fontweight="bold", color="#7C3AED")

    # 前端
    rect = FancyBboxPatch((9.5, 3.5), 1.3, 1.0, boxstyle="round,pad=0.1",
                          facecolor=(0.86, 0.98, 0.87), edgecolor="#94A3B8", linewidth=1)
    ax.add_patch(rect)
    ax.text(10.15, 4.0, "前端\nHTML/JS", ha="center", va="center", **font_prop)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="#F7F7F8")
    plt.close(fig)
    buf.seek(0)
    slide.shapes.add_picture(buf, Emu(MARGIN), Emu(inches(1.5)), Emu(CONTENT_W), Emu(inches(5.0)))

    # ========== 面试考点 1-5 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.4), CONTENT_W, inches(0.7),
             "面试考点 (1/2)", size=28, bold=True, color=ACCENT)
    qa_pairs_1 = [
        ("Q1: 什么是 RAG？和微调有什么区别？",
         "A: RAG = 检索增强生成，先从文档库检索相关内容再让 LLM 基于检索结果回答。"
         "微调是修改模型参数让它学会新知识。RAG 便宜、实时更新、不需要训练；"
         "微调效果更深入但贵、慢、需要大量数据。"),

        ("Q2: RAG 的完整流程是什么？",
         "A: 离线阶段: 文档 -> 分块(Chunking) -> 向量化(Embedding) -> 存入向量库。"
         "在线阶段: 用户提问 -> 向量化 -> 向量库检索 top-k -> 拼接到 prompt -> LLM 生成回答。"),

        ("Q3: Embedding 是什么？为什么能做语义搜索？",
         "A: Embedding 把文本转成高维向量（一组数字），语义相近的文本在向量空间中距离近。"
         "传统搜索靠关键词匹配，向量搜索靠语义匹配，所以 'React 组件' 能搜到 'Vue 组件'。"),

        ("Q4: chunk_size 怎么选？太大太小有什么问题？",
         "A: 太大 -> 噪音多，检索不精准，浪费 token。太小 -> 丢失上下文，语义不完整。"
         "通常 300-500 字符起步，用评估集实验对比。overlap 50-100 字防止切断关键信息。"),

        ("Q5: 什么是 Rerank？什么时候需要？",
         "A: 向量检索（Bi-encoder）是粗排，Rerank（Cross-encoder）是精排。"
         "Bi-encoder 快但精度一般，Cross-encoder 慢但精准。"
         "流程: Bi-encoder 召回 20 条 -> Cross-encoder 重新打分取 3 条。"
         "文档量大(>1000条)且精度不够时需要 Rerank。"),
    ]
    add_qa_list(slide, MARGIN, inches(1.2), CONTENT_W, inches(5.8), qa_pairs_1)

    # ========== 面试考点 6-10 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text(slide, MARGIN, inches(0.4), CONTENT_W, inches(0.7),
             "面试考点 (2/2)", size=28, bold=True, color=ACCENT)
    qa_pairs_2 = [
        ("Q6: 怎么评估 RAG 系统的效果？",
         "A: 分开评估: (1)检索评估 — 期望来源是否在 top-k 结果中（命中率）; "
         "(2)生成评估 — AI 回答是否包含期望关键词（通过率）。"
         "检索差 -> 优化分块/embedding/rerank; 生成差 -> 优化 prompt/top_k。"),

        ("Q7: FastAPI 和 Express 的主要区别？",
         "A: FastAPI 自带 Pydantic 类型校验（Express 需要 Zod）, "
         "自动生成 API 文档 /docs（Express 需要手动配 Swagger）, "
         "原生 async/await 支持。路由用装饰器 @app.get() 而不是 app.get(path, handler)。"),

        ("Q8: SSE 和 WebSocket 有什么区别？AI 场景用哪个？",
         "A: SSE 是单向（服务器 -> 客户端），基于 HTTP，简单。"
         "WebSocket 是双向，需要额外握手，适合聊天室/游戏。"
         "AI 流式输出用 SSE 就够了（ChatGPT 就是用的 SSE），不需要 WebSocket。"),

        ("Q9: 前端怎么接收 SSE 流？为什么不用 EventSource？",
         "A: 因为 EventSource 只支持 GET，RAG 问答需要 POST 发送请求体。"
         "所以用 fetch + ReadableStream: response.body.getReader() 逐块读取，"
         "TextDecoder 解码二进制，按 SSE 格式解析 event/data 字段。"),

        ("Q10: RAG 系统用户问了文档中没有的问题怎么办？",
         "A: 设置相似度阈值（如 L2 距离 > 1.5 视为不相关），低于阈值则拒绝回答，"
         "返回'文档中未找到相关信息'。同时在 prompt 中明确要求 LLM: "
         "'如果文档中没有相关信息，明确说无法回答'，双重保障。"),
    ]
    add_qa_list(slide, MARGIN, inches(1.2), CONTENT_W, inches(5.8), qa_pairs_2)

    # ========== 总结页 ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, ACCENT)
    add_text(slide, MARGIN, inches(2.0), CONTENT_W, inches(1.0),
             "Week 3 Complete!", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN, inches(3.2), CONTENT_W, inches(0.6),
             "RAG 从概念到全栈落地", size=20, color=WHITE, align=PP_ALIGN.CENTER)
    summary_items = [
        "RAG: 分块 -> Embedding -> 向量检索 -> LLM 生成",
        "优化: chunk_size / top_k / 阈值 / Rerank / 评估",
        "FastAPI: 路由 + Pydantic + SSE + CORS",
        "全栈: 后端 API + 前端 SSE 流式聊天界面",
    ]
    add_bullet_list(slide, inches(3.0), inches(4.2), inches(7.333), inches(2.5),
                    summary_items, size=16, color=RGBColor(0xBF, 0xDB, 0xFE), spacing=Pt(10))

    # ========== 保存 ==========
    output = Path(__file__).parent / "week3_summary.pptx"
    prs.save(output)
    print(f"PPT 已保存到: {output}")
    return output


if __name__ == "__main__":
    generate_ppt()
