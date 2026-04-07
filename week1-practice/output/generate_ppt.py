"""
生成第一周 Python 学习汇总 PPT
运行: python week1-practice/output/generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path


# ===========================================
# 颜色主题
# ===========================================
DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)       # 深色背景
LIGHT_BG = RGBColor(0x28, 0x28, 0x3E)       # 浅深色背景
ACCENT = RGBColor(0x89, 0xB4, 0xFA)         # 蓝色强调
ACCENT2 = RGBColor(0xA6, 0xE3, 0xA1)        # 绿色
ACCENT3 = RGBColor(0xFA, 0xB3, 0x87)        # 橙色
WHITE = RGBColor(0xF5, 0xF5, 0xF5)          # 白色文字
GRAY = RGBColor(0xA0, 0xA0, 0xB0)           # 灰色副标题
YELLOW = RGBColor(0xF9, 0xE2, 0xAF)         # 黄色


def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_text(slide, text, left, top, width, height, font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_bullet_list(slide, items, left, top, width, height, font_size=16, color=WHITE, bullet_color=ACCENT):
    """添加带样式的列表"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # 用彩色符号作为 bullet
        run1 = p.add_run()
        run1.text = "  >  "
        run1.font.size = Pt(font_size)
        run1.font.color.rgb = bullet_color
        run1.font.bold = True

        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color

        p.space_after = Pt(6)

    return tf


def add_code_block(slide, code, left, top, width, height, font_size=13):
    """添加代码块"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    # 背景用 shape 模拟
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(left - 0.1), Inches(top - 0.1),
        Inches(width + 0.2), Inches(height + 0.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x18, 0x18, 0x28)
    shape.line.fill.background()

    # 把代码框放到前面
    # (shape 先添加所以在下层，文字在上层)

    for i, line in enumerate(code.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT2
        p.font.name = "Consolas"
        p.space_after = Pt(2)

    return tf


def add_two_column_compare(slide, left_title, left_items, right_title, right_items, top=2.0):
    """添加两列对比"""
    # 左标题
    add_title_text(slide, left_title, 0.5, top, 4, 0.5, font_size=20, color=ACCENT3)
    add_bullet_list(slide, left_items, 0.5, top + 0.5, 4.2, 3, font_size=14, bullet_color=ACCENT3)

    # 右标题
    add_title_text(slide, right_title, 5.2, top, 4, 0.5, font_size=20, color=ACCENT2)
    add_bullet_list(slide, right_items, 5.2, top + 0.5, 4.5, 3, font_size=14, bullet_color=ACCENT2)


def add_tag(slide, text, left, top, color=ACCENT):
    """添加小标签"""
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(len(text) * 0.13 + 0.3), Inches(0.35)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


# ===========================================
# 生成 PPT
# ===========================================

def generate_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # ==========================================
    # 第 1 页：封面
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "Python 第一周学习总结", 0.8, 1.2, 8, 1, font_size=40, color=WHITE)
    add_title_text(slide, "From JavaScript to Python + AI", 0.8, 2.2, 8, 0.6, font_size=22, color=ACCENT, bold=False)
    add_title_text(slide, "2026 年 4 月  |  前端工程师的 Python 之旅", 0.8, 3.2, 8, 0.5, font_size=16, color=GRAY, bold=False)

    # ==========================================
    # 第 2 页：学习路线图
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "学习路线图", 0.5, 0.3, 9, 0.6, font_size=32, color=ACCENT)

    days = [
        ("Day 1", "Python 基础", "变量、字符串、列表、字典、函数、循环", ACCENT),
        ("Day 2", "进阶语法", "类、继承、模块、异常处理、异步", ACCENT2),
        ("Day 3", "环境与 API", "venv、pip、.env、API 调用", ACCENT3),
        ("Day 4", "LLM 概念", "Token、Temperature、Prompt、Few-shot", YELLOW),
        ("Day 5", "聊天机器人", "模块架构、流式输出、命令系统", ACCENT),
        ("Day 6", "综合项目", "笔记生成器、异步并发、装饰器", ACCENT2),
    ]

    for i, (day, title, desc, color) in enumerate(days):
        y = 1.2 + i * 0.65
        add_tag(slide, day, 0.5, y, color)
        add_title_text(slide, title, 1.5, y - 0.02, 2.5, 0.4, font_size=18, color=WHITE)
        add_title_text(slide, desc, 4.0, y - 0.02, 5.5, 0.4, font_size=14, color=GRAY, bold=False)

    # ==========================================
    # 第 3 页：Day 1 Python 基础
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "Day 1: Python 基础", 0.5, 0.3, 9, 0.6, font_size=30, color=ACCENT)
    add_title_text(slide, "对比 JavaScript 快速上手", 0.5, 0.85, 9, 0.4, font_size=16, color=GRAY, bold=False)

    add_two_column_compare(slide,
        "JavaScript", [
            "const name = 'Alice'",
            "`Hello ${name}`",
            "arr.push() / arr.map()",
            "obj.key / obj['key']",
            "function fn() { }",
            "arr.slice(1, 4)",
        ],
        "Python", [
            "name = 'Alice'  (无需声明)",
            "f\"Hello {name}\"",
            "list.append() / 列表推导式",
            "dict['key'] / dict.get('key')",
            "def fn():",
            "arr[1:4]  (切片语法)",
        ],
        top=1.5
    )

    # ==========================================
    # 第 4 页：Day 2 进阶
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "Day 2: 进阶语法", 0.5, 0.3, 9, 0.6, font_size=30, color=ACCENT2)

    add_bullet_list(slide, [
        "类: class + __init__() + self（JS: constructor + this）",
        "继承: class Child(Parent)（JS: class Child extends Parent）",
        "模块: from module import x（JS: import x from 'module'）",
        "异常: try / except / finally（JS: try / catch / finally）",
        "异步: async def + await（JS: async function + await）",
        "类型注解: x: int = 5（类似 TypeScript，但不强制）",
    ], 0.5, 1.2, 9, 3.5, font_size=16)

    add_code_block(slide, (
        'class Task:\n'
        '    def __init__(self, title):  # self = JS 的 this\n'
        '        self.title = title\n'
        '    def complete(self):\n'
        '        self.done = True'
    ), 5.5, 3.5, 4, 1.5, font_size=12)

    # ==========================================
    # 第 5 页：Day 3 环境与 API
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "Day 3: 环境管理与 API 调用", 0.5, 0.3, 9, 0.6, font_size=30, color=ACCENT3)

    add_two_column_compare(slide,
        "npm (JS)", [
            "npm init",
            "npm install axios",
            "package.json",
            ".env + dotenv",
            "node index.js",
        ],
        "pip (Python)", [
            "python -m venv venv",
            "pip install anthropic",
            "requirements.txt",
            ".env + python-dotenv",
            "python main.py",
        ],
        top=1.3
    )

    add_code_block(slide, (
        'from anthropic import Client\n'
        'client = Client(api_key=os.getenv("API_KEY"))\n'
        'response = client.messages.create(\n'
        '    model="...", max_tokens=200,\n'
        '    messages=[{"role": "user", "content": "Hi"}]\n'
        ')'
    ), 0.5, 3.5, 9, 1.5, font_size=12)

    # ==========================================
    # 第 6 页：Day 4 LLM 核心概念
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "Day 4: LLM 核心概念", 0.5, 0.3, 9, 0.6, font_size=30, color=YELLOW)

    concepts = [
        ("Token", "文本的最小单位（英文 ~1 word/token，中文 ~1-2 字/token）"),
        ("Temperature", "随机性控制：0=确定性（代码），0.7-1.0=创意写作"),
        ("System Prompt", "设定 AI 角色和行为约束"),
        ("Few-shot", "给 2-3 个示例让 AI 照格式回答，比描述格式更有效"),
        ("结构化输出", "通过 Prompt 约束 AI 返回 JSON 格式"),
    ]

    for i, (term, desc) in enumerate(concepts):
        y = 1.2 + i * 0.7
        add_tag(slide, term, 0.5, y, YELLOW)
        width = len(term) * 0.13 + 0.3
        add_title_text(slide, desc, 0.5 + width + 0.3, y + 0.02, 8 - width, 0.4, font_size=15, color=WHITE, bold=False)

    # ==========================================
    # 第 7 页：Day 5 聊天机器人
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "Day 5: CLI 聊天机器人项目", 0.5, 0.3, 9, 0.6, font_size=30, color=ACCENT)

    add_bullet_list(slide, [
        "模块拆分: config.py / client.py / commands.py",
        "@dataclass 管理配置（类似 TS interface）",
        "多轮对话: 维护 message history 传给 API",
        "流式输出: 逐字打印，提升用户体验",
        "命令系统: /help /clear /save /history",
        "持久化: JSON 保存/加载对话记录",
    ], 0.5, 1.1, 5, 3.5, font_size=15)

    add_code_block(slide, (
        'chatbot/\n'
        '  __init__.py    # 包入口\n'
        '  config.py      # @dataclass 配置\n'
        '  client.py      # ChatClient 类\n'
        '  commands.py    # 命令处理器'
    ), 5.8, 1.3, 3.8, 1.8, font_size=13)

    # ==========================================
    # 第 8 页：Day 6 综合项目
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "Day 6: 综合项目 — 代码笔记生成器", 0.5, 0.3, 9, 0.6, font_size=28, color=ACCENT2)

    add_bullet_list(slide, [
        "扫描目录下所有 .py 文件（Path.rglob）",
        "异步并发调用 AI（asyncio.gather = Promise.all）",
        "AI 返回 JSON 结构化学习笔记",
        "@timer 装饰器计时 + @retry 重试机制",
        "输出 Markdown 报告 + JSON 原始数据",
    ], 0.5, 1.1, 5, 3, font_size=15)

    add_code_block(slide, (
        'note_generator/\n'
        '  __init__.py     # 统一导出\n'
        '  config.py       # 环境变量管理\n'
        '  scanner.py      # 文件扫描器\n'
        '  generator.py    # AI 笔记生成\n'
        '\n'
        'output/\n'
        '  notes_report.md # Markdown 报告\n'
        '  notes_data.json # JSON 数据'
    ), 5.8, 1.1, 3.8, 2.8, font_size=13)

    # 异步并发对比
    add_code_block(slide, (
        '# 串行（慢）          # 并行（快）\n'
        'for f in files:      tasks = [gen(f) for f in files]\n'
        '  await gen(f)       await asyncio.gather(*tasks)'
    ), 0.5, 3.7, 9, 1.2, font_size=12)

    # ==========================================
    # 第 9 页：Python vs JS 对照表
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "Python vs JavaScript 核心对照", 0.5, 0.2, 9, 0.6, font_size=30, color=ACCENT)

    comparisons = [
        ("变量声明", "const / let / var", "直接赋值，无需关键字"),
        ("布尔值", "true / false", "True / False（大写）"),
        ("空值", "null / undefined", "None（只有一个）"),
        ("字符串", "`${x}` 模板字符串", 'f"{x}" f-string'),
        ("数组/列表", "[].map() .filter()", "列表推导式 [x for x in ...]"),
        ("对象/字典", "obj.key", "dict['key'] / dict.get('key')"),
        ("类构造", "constructor()", "__init__(self)"),
        ("this/self", "this（隐式）", "self（显式，必须写）"),
        ("模块导入", "import x from 'y'", "from y import x"),
        ("错误处理", "try / catch", "try / except"),
        ("代码块", "{ } 花括号", "缩进（4空格）"),
    ]

    # 表头
    add_title_text(slide, "概念", 0.3, 0.85, 1.5, 0.35, font_size=13, color=DARK_BG, bold=True)
    add_title_text(slide, "JavaScript", 2.0, 0.85, 3.5, 0.35, font_size=13, color=DARK_BG, bold=True)
    add_title_text(slide, "Python", 5.5, 0.85, 4.2, 0.35, font_size=13, color=DARK_BG, bold=True)

    # 表头背景
    header_bg = slide.shapes.add_shape(1, Inches(0.2), Inches(0.82), Inches(9.5), Inches(0.38))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = ACCENT
    header_bg.line.fill.background()

    for i, (concept, js, py) in enumerate(comparisons):
        y = 1.25 + i * 0.36

        # 交替行背景
        if i % 2 == 0:
            row_bg = slide.shapes.add_shape(1, Inches(0.2), Inches(y - 0.03), Inches(9.5), Inches(0.36))
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = LIGHT_BG
            row_bg.line.fill.background()

        add_title_text(slide, concept, 0.3, y, 1.5, 0.3, font_size=12, color=YELLOW, bold=True)
        add_title_text(slide, js, 2.0, y, 3.5, 0.3, font_size=12, color=ACCENT3, bold=False)
        add_title_text(slide, py, 5.5, y, 4.2, 0.3, font_size=12, color=ACCENT2, bold=False)

    # ==========================================
    # 第 10 页：练习完成情况
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "练习完成情况", 0.5, 0.3, 9, 0.6, font_size=30, color=ACCENT)

    exercises = [
        ("practice_01", "Python 基础", "DONE", "列表操作、字典处理、成绩计算"),
        ("practice_02", "类与模块", "DONE", "Task 类、TaskManager、继承"),
        ("practice_03", "文件 I/O", "DONE", "JSON 读写、文件扫描、报告生成"),
        ("practice_04", "API 调用", "DONE", "角色扮演、流式输出、retry 装饰器、few-shot"),
        ("practice_05", "综合项目", "DONE", "笔记生成器（异步并发 + 模块拆分）"),
    ]

    for i, (name, topic, status, detail) in enumerate(exercises):
        y = 1.2 + i * 0.75
        status_color = ACCENT2

        add_tag(slide, status, 0.5, y, status_color)
        add_title_text(slide, f"{name}: {topic}", 1.5, y - 0.02, 3, 0.4, font_size=17, color=WHITE)
        add_title_text(slide, detail, 1.5, y + 0.32, 8, 0.3, font_size=13, color=GRAY, bold=False)

    # ==========================================
    # 第 11 页：参考文件
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "自建参考手册", 0.5, 0.3, 9, 0.6, font_size=30, color=ACCENT3)

    refs = [
        ("reference_list_dict.py", "列表和字典常用操作",
         "append/pop/sort/sorted、lambda、列表推导式、\ndict 增删查改、enumerate/zip/any/all/min/max/sum"),
        ("reference_json.py", "JSON 操作大全",
         "dumps/loads/dump/load、嵌套操作、\n自定义序列化(default=str)、深拷贝技巧"),
        ("07_new_syntax.py", "Python 3.8~3.13 新语法",
         ":= 海象运算符、| 字典合并、match/case 模式匹配、\nexcept* 异常组、f-string 增强"),
    ]

    for i, (filename, title, content) in enumerate(refs):
        y = 1.1 + i * 1.3
        add_tag(slide, filename, 0.5, y, ACCENT3)
        add_title_text(slide, title, 0.5, y + 0.45, 9, 0.4, font_size=17, color=WHITE)
        add_title_text(slide, content, 0.5, y + 0.85, 9, 0.6, font_size=13, color=GRAY, bold=False)

    # ==========================================
    # 第 12 页：Python 新语法速查
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "Python 新语法速查（你是 3.11）", 0.5, 0.3, 9, 0.6, font_size=28, color=ACCENT)

    syntax_items = [
        ("3.8", ":= 海象运算符", "if (n := len(a)) > 10:", "DONE"),
        ("3.9", "| 字典合并", "merged = dict_a | dict_b", "DONE"),
        ("3.10", "match/case", "match status: case 200: ...", "DONE"),
        ("3.10", "int | str 类型注解", "def f(x: int | str):", "DONE"),
        ("3.11", "except* 异常组", "except* ValueError as eg:", "DONE"),
        ("3.12", "f-string 嵌套引号", 'f"name: {d["key"]}"', "N/A"),
        ("3.13", "JIT 编译器", "自动优化热点代码", "N/A"),
    ]

    for i, (ver, feature, example, status) in enumerate(syntax_items):
        y = 1.1 + i * 0.58
        color = ACCENT2 if status == "DONE" else GRAY

        add_tag(slide, ver, 0.5, y, color)
        add_title_text(slide, feature, 1.3, y - 0.02, 2.5, 0.35, font_size=14, color=WHITE)
        add_title_text(slide, example, 4.0, y - 0.02, 4.5, 0.35, font_size=13, color=color, bold=False)

        status_text = "CAN USE" if status == "DONE" else "3.12+"
        status_c = ACCENT2 if status == "DONE" else RGBColor(0x60, 0x60, 0x70)
        add_title_text(slide, status_text, 8.8, y - 0.02, 1, 0.35, font_size=11, color=status_c)

    # ==========================================
    # 第 13 页：关键学习心得
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "10 个关键心得", 0.5, 0.2, 9, 0.5, font_size=30, color=YELLOW)

    tips = [
        "缩进即语法 — 4 空格不是风格，是语法要求",
        "self 必须显式写 — 不像 JS 的 this 是隐式的",
        "列表推导式 > filter + lambda — 更 Pythonic",
        "dict.get() > dict[] — 安全取值不报错",
        "with open() — 自动关闭文件，防止资源泄漏",
        "f-string — 最好用的字符串格式化",
        "Path > os.path — 更现代，可以链式操作",
        "生成器省内存 — sum(1 for ...) 比 len(list) 高效",
        "asyncio.gather = Promise.all — 并行调用的关键",
        "JSON 是字符串，dict 是数据结构 — 别混淆",
    ]

    for i, tip in enumerate(tips):
        y = 0.8 + i * 0.44
        num_color = ACCENT if i < 5 else ACCENT2
        add_title_text(slide, f" {i+1}.", 0.3, y, 0.5, 0.35, font_size=14, color=num_color)
        add_title_text(slide, tip, 0.9, y, 8.5, 0.35, font_size=14, color=WHITE, bold=False)

    # ==========================================
    # 第 14 页：项目架构
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "综合项目架构", 0.5, 0.3, 9, 0.6, font_size=30, color=ACCENT2)

    # 流程图用文本模拟
    add_code_block(slide, (
        '  [Scanner]              [NoteGenerator]              [Output]\n'
        '      |                        |                         |\n'
        ' Path.rglob()      AsyncAnthropic + gather()      Markdown + JSON\n'
        '      |                        |                         |\n'
        ' .py files  ------>   AI parallel calls  ------>   notes_report.md\n'
        '                                                   notes_data.json'
    ), 0.5, 1.1, 9, 2, font_size=13)

    add_bullet_list(slide, [
        "Config: 集中管理环境变量（.env + 验证）",
        "Scanner: 递归扫描 + 文件元数据（行数/大小）",
        "NoteGenerator: 异步 AI 调用 + JSON 解析 + 错误兜底",
        "@timer 装饰器: 同时支持 sync / async 函数",
        "命令行参数: sys.argv 指定扫描目录",
    ], 0.5, 3.3, 9, 2, font_size=14)

    # ==========================================
    # 第 15 页：总结
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_title_text(slide, "Week 1 Complete!", 0.8, 1.0, 8, 1, font_size=44, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_title_text(slide, "从 JavaScript 到 Python + AI", 0.8, 2.0, 8, 0.6, font_size=22, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

    stats = "6 天课程  |  5 个练习  |  3 份参考手册  |  1 个综合项目"
    add_title_text(slide, stats, 0.8, 2.8, 8, 0.5, font_size=16, color=GRAY, bold=False, alignment=PP_ALIGN.CENTER)

    add_title_text(slide, "Next: Week 2 >>>", 0.8, 3.8, 8, 0.5, font_size=20, color=ACCENT2, alignment=PP_ALIGN.CENTER)

    # ==========================================
    # 保存
    # ==========================================
    output_path = Path(__file__).parent / "week1_summary.pptx"
    prs.save(str(output_path))
    print(f"PPT saved: {output_path}")
    return output_path


if __name__ == "__main__":
    generate_ppt()